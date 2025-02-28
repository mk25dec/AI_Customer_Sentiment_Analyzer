import os
import win32com.client
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
from PIL import Image

# File paths
BASE_PATH = r"C:\Users\ManojChauhan\OneDrive - Radicant\Documents\Projects\charts_auto"
EXCEL_FILE_PATH = os.path.join(BASE_PATH, "charts.xlsx")
SHEET_NAME = "NewChart"
PPTX_FILE_PATH = os.path.join(BASE_PATH, "charts_presentation.pptx")
IMAGE_FOLDER = os.path.join(BASE_PATH, "chart_images")

# Generate timestamp (YYYYMMDDHHMMSS)
TIMESTAMP = datetime.now().strftime("%Y%m%d%H%M%S")

# Chart configuration: slide number and position (in cm)
CHART_MAPPING = {
    "SalesChart": {"slide_number": 1, "left": 2, "top": 2, "width": 15, "height": 10},
    "ProfitChart": {"slide_number": 3, "left": 4, "top": 4, "width": 12, "height": 8},
    "GrowthChart": {"slide_number": 5, "left": 3, "top": 3, "width": 14, "height": 9},
}


# 🔹 1. Archive Existing Files (Rename Instead of Delete)
def archive_existing_files():
    """Renames old files and folders instead of deleting them."""

    def rename_file(path):
        """Renames a file by inserting the timestamp before its extension."""
        if os.path.exists(path):
            base, ext = os.path.splitext(path)  # Split filename and extension
            new_path = f"{base}_{TIMESTAMP}{ext}"  # Insert timestamp before extension
            try:
                os.rename(path, new_path)
                print(f"Renamed: {path} → {new_path}")
            except Exception as e:
                print(f"⚠️ Could not rename {path}: {e}")

    def rename_folder(path):
        """Renames a folder by appending the timestamp."""
        if os.path.exists(path):
            new_path = f"{path}_{TIMESTAMP}"
            try:
                os.rename(path, new_path)
                print(f"Renamed: {path} → {new_path}")
            except Exception as e:
                print(f"⚠️ Could not rename {path}: {e}")

    # Archive the PowerPoint file
    rename_file(PPTX_FILE_PATH)
    # Archive the chart images folder
    rename_folder(IMAGE_FOLDER)


# 🔹 2. Create New Folder and Files
def setup_environment():
    """Creates a fresh chart images folder and a new PowerPoint file."""
    os.makedirs(IMAGE_FOLDER, exist_ok=True)
    print(f"✅ Created new chart images folder: {IMAGE_FOLDER}")

    prs = Presentation()
    prs.save(PPTX_FILE_PATH)
    print(f"✅ Created a new PowerPoint file: {PPTX_FILE_PATH}")


# 🔹 3. Extract and Save Charts from Excel
def extract_charts_from_excel():
    """Extracts charts from the Excel file and saves them as images."""
    excel = win32com.client.Dispatch("Excel.Application")
    excel.Visible = False
    wb = excel.Workbooks.Open(EXCEL_FILE_PATH)
    ws = wb.Sheets(SHEET_NAME)

    chart_images = {}

    for chart_obj in ws.ChartObjects():
        chart = chart_obj.Chart
        chart_name = chart_obj.Name.strip()  # Extract chart name & remove spaces

        if chart_name in CHART_MAPPING:
            image_path = os.path.join(IMAGE_FOLDER, f"{chart_name}.png")
            chart.Export(image_path)
            chart_images[chart_name] = image_path

    wb.Close(False)
    excel.Quit()

    return chart_images


# 🔹 4. Create PowerPoint and Insert Charts
def insert_charts_into_ppt(chart_images):
    """Inserts valid chart images into the PowerPoint file at specified slides."""
    prs = Presentation(PPTX_FILE_PATH)

    # Function to verify if an image is valid
    def is_valid_image(image_path):
        try:
            with Image.open(image_path) as img:
                img.verify()  # Only verifies, does not load full image
            return True
        except Exception as e:
            print(f"⚠️ Invalid image detected: {image_path} | Error: {e}")
            return False

    # Filter out invalid images
    valid_chart_images = {name: path for name, path in chart_images.items() if is_valid_image(path)}

    # Ensure PowerPoint has enough slides
    if valid_chart_images:
        max_slide_number = max(CHART_MAPPING[c]["slide_number"] for c in valid_chart_images if c in CHART_MAPPING)
        while len(prs.slides) < max_slide_number:
            prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    else:
        print("⚠️ No valid charts found! Please check chart names in Excel and `CHART_MAPPING`.")
        return

    # Function to convert cm to inches
    def cm_to_inches(cm):
        return cm / 2.54

    # Insert images at specified slides
    for chart_name, image_path in valid_chart_images.items():
        config = CHART_MAPPING[chart_name]
        slide_index = config["slide_number"] - 1  # Convert to zero-based index
        slide = prs.slides[slide_index]

        left = Inches(cm_to_inches(config["left"]))
        top = Inches(cm_to_inches(config["top"]))
        width = Inches(cm_to_inches(config["width"]))
        height = Inches(cm_to_inches(config["height"]))

        slide.shapes.add_picture(image_path, left, top, width=width, height=height)

    prs.save(PPTX_FILE_PATH)

    print(f"✅ Inserted {len(valid_chart_images)} charts into specific slides in {PPTX_FILE_PATH}.")


# 🔹 Main Execution Flow
def main():
    """Runs the full process in a structured manner."""
    archive_existing_files()  # Step 1: Rename old files instead of deleting
    setup_environment()  # Step 2: Create fresh folders & new PowerPoint
    chart_images = extract_charts_from_excel()  # Step 3: Save charts from Excel
    insert_charts_into_ppt(chart_images)  # Step 4: Insert charts into PowerPoint


# Run the script
if __name__ == "__main__":
    main()
